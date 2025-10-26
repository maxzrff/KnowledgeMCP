"""
PPTX document processor.
"""

from pathlib import Path
from typing import Any

from pptx import Presentation

from src.models.document import DocumentFormat
from src.processors.base import BaseProcessor
from src.utils.logging_config import get_logger

logger = get_logger(__name__)


class PPTXProcessor(BaseProcessor):
    """PPTX document processor."""

    @property
    def supported_format(self) -> DocumentFormat:
        return DocumentFormat.PPTX

    async def extract_text(self, file_path: Path) -> str:
        """Extract text from PPTX."""
        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

            text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(text)} characters from PPTX: {file_path.name}")
            return text
        except Exception as e:
            logger.error(f"Failed to extract text from PPTX {file_path}: {e}")
            raise

    async def extract_metadata(self, file_path: Path) -> dict[str, Any]:
        """Extract metadata from PPTX."""
        try:
            prs = Presentation(file_path)
            metadata = {
                "slide_count": len(prs.slides),
                "format": "pptx",
            }

            # Core properties
            if prs.core_properties.author:
                metadata["author"] = prs.core_properties.author
            if prs.core_properties.title:
                metadata["title"] = prs.core_properties.title

            return metadata
        except Exception as e:
            logger.warning(f"Failed to extract PPTX metadata: {e}")
            return {"format": "pptx"}
